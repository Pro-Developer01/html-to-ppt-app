#!/usr/bin/env python3
"""
Script to convert HTML slides to PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation with 16:9 aspect ratio (1280x720)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Slide 1: Cover/Proposal Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background color
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(243, 244, 253)  # #F3F4FD
    
    # Brand logo
    logo_box = slide1.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(3), Inches(0.5))
    logo_frame = logo_box.text_frame
    logo_frame.text = "CodeXBit"
    logo_para = logo_frame.paragraphs[0]
    logo_run = logo_para.runs[0]
    logo_run.font.size = Pt(32)
    logo_run.font.bold = True
    logo_run.font.color.rgb = RGBColor(30, 30, 63)
    logo_run.font.name = 'Poppins'
    
    # Proposal Document badge (top right)
    badge_box = slide1.shapes.add_textbox(Inches(7.5), Inches(0.4), Inches(2), Inches(0.4))
    badge_frame = badge_box.text_frame
    badge_frame.text = "Proposal Document"
    badge_para = badge_frame.paragraphs[0]
    badge_para.alignment = PP_ALIGN.CENTER
    badge_run = badge_para.runs[0]
    badge_run.font.size = Pt(11)
    badge_run.font.color.rgb = RGBColor(255, 255, 255)
    badge_run.font.name = 'Inter'
    # Add shape fill
    badge_shape = badge_box
    badge_shape.fill.solid()
    badge_shape.fill.fore_color.rgb = RGBColor(30, 30, 63)
    
    # Main headline
    title_box = slide1.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Proposal"
    title_para = title_frame.paragraphs[0]
    title_run = title_para.runs[0]
    title_run.font.size = Pt(64)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(30, 30, 63)
    title_run.font.name = 'Poppins'
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(5), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Presented by #TeamCodeXBit"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_run = subtitle_para.runs[0]
    subtitle_run.font.size = Pt(20)
    subtitle_run.font.color.rgb = RGBColor(75, 85, 99)
    subtitle_run.font.name = 'Inter'
    
    # Feature pills
    pills = [
        ("Scalable Solutions", 0.6),
        ("Secure", 2.8),
        ("High Performance", 4.3)
    ]
    
    for pill_text, x_pos in pills:
        pill_box = slide1.shapes.add_textbox(Inches(x_pos), Inches(3.5), Inches(1.5), Inches(0.4))
        pill_frame = pill_box.text_frame
        pill_frame.text = pill_text
        pill_para = pill_frame.paragraphs[0]
        pill_para.alignment = PP_ALIGN.CENTER
        pill_run = pill_para.runs[0]
        pill_run.font.size = Pt(11)
        pill_run.font.color.rgb = RGBColor(107, 92, 231)
        pill_run.font.bold = True
        pill_run.font.name = 'Inter'
        # Add shape fill
        pill_box.fill.solid()
        pill_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Right side card with icons
    card_box = slide1.shapes.add_textbox(Inches(6.2), Inches(1.5), Inches(3.2), Inches(3))
    card_frame = card_box.text_frame
    card_frame.text = "Engineering Intelligence\n\n• Web Dev\n• Mobile App\n• Cloud\n• Automation"
    card_para = card_frame.paragraphs[0]
    card_run = card_para.runs[0]
    card_run.font.size = Pt(14)
    card_run.font.bold = True
    card_run.font.color.rgb = RGBColor(30, 30, 63)
    card_run.font.name = 'Inter'
    # Style the card
    card_box.fill.solid()
    card_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Date
    date_box = slide1.shapes.add_textbox(Inches(8.5), Inches(5), Inches(1), Inches(0.3))
    date_frame = date_box.text_frame
    date_frame.text = "2026"
    date_para = date_frame.paragraphs[0]
    date_para.alignment = PP_ALIGN.RIGHT
    date_run = date_para.runs[0]
    date_run.font.size = Pt(11)
    date_run.font.color.rgb = RGBColor(107, 114, 128)
    date_run.font.name = 'Inter'
    
    # Slide 2: About CodeXBit
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background2 = slide2.background
    fill2 = background2.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(243, 244, 253)
    
    # Section tag
    tag_box = slide2.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(3), Inches(0.3))
    tag_frame = tag_box.text_frame
    tag_frame.text = "WHO WE ARE"
    tag_para = tag_frame.paragraphs[0]
    tag_run = tag_para.runs[0]
    tag_run.font.size = Pt(11)
    tag_run.font.bold = True
    tag_run.font.color.rgb = RGBColor(107, 92, 231)
    tag_run.font.name = 'Inter'
    
    # Main title
    title2_box = slide2.shapes.add_textbox(Inches(0.6), Inches(1), Inches(5), Inches(0.6))
    title2_frame = title2_box.text_frame
    title2_frame.text = "About CodeXBit"
    title2_para = title2_frame.paragraphs[0]
    title2_run = title2_para.runs[0]
    title2_run.font.size = Pt(40)
    title2_run.font.bold = True
    title2_run.font.color.rgb = RGBColor(30, 30, 63)
    title2_run.font.name = 'Poppins'
    
    # Subtitle
    subtitle2_box = slide2.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(5), Inches(0.4))
    subtitle2_frame = subtitle2_box.text_frame
    subtitle2_frame.text = "A Company Leading By Young Engineers, Entrepreneurs and Innovative Team"
    subtitle2_para = subtitle2_frame.paragraphs[0]
    subtitle2_run = subtitle2_para.runs[0]
    subtitle2_run.font.size = Pt(13)
    subtitle2_run.font.color.rgb = RGBColor(75, 85, 99)
    subtitle2_run.font.name = 'Inter'
    
    # Content paragraphs
    content_text = """Founded in 2024, CodexBit is a modern product and technology company focused on building scalable, secure, and high-performance digital solutions. What began as a small engineering team has grown into a full-service IT company delivering reliable software for businesses worldwide.

Over the past two years, we have successfully delivered 50+ projects for startups, SMEs, and fast-growing enterprises. Beyond client services, CodexBit actively designs, builds, and operates its own products, including BoomGhoom and ConnectCRM, reflecting our strong product mindset and long-term vision."""
    
    content_box = slide2.shapes.add_textbox(Inches(0.6), Inches(2.3), Inches(4.5), Inches(2))
    content_frame = content_box.text_frame
    content_frame.text = content_text
    content_frame.word_wrap = True
    for para in content_frame.paragraphs:
        para.font.size = Pt(11)
        para.font.color.rgb = RGBColor(75, 85, 99)
        para.font.name = 'Inter'
        para.line_spacing = 1.4
    
    # Badge
    exp_badge = slide2.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(3), Inches(0.4))
    exp_frame = exp_badge.text_frame
    exp_frame.text = "10+ Years of Experienced Team"
    exp_para = exp_frame.paragraphs[0]
    exp_para.alignment = PP_ALIGN.CENTER
    exp_run = exp_para.runs[0]
    exp_run.font.size = Pt(11)
    exp_run.font.bold = True
    exp_run.font.color.rgb = RGBColor(79, 70, 229)
    exp_run.font.name = 'Inter'
    exp_badge.fill.solid()
    exp_badge.fill.fore_color.rgb = RGBColor(238, 242, 255)
    
    # Stats section (right side)
    stats = [
        ("50+", "Successful Projects", 6.2, 1.2),
        ("20+", "Clients", 6.2, 2.5),
        ("35+", "Team Members", 7.8, 2.5),
        ("2+", "Products", 6.2, 3.8),
        ("950+", "Features Built", 7.8, 3.8)
    ]
    
    for stat_num, stat_label, x_pos, y_pos in stats:
        stat_box = slide2.shapes.add_textbox(Inches(x_pos), Inches(y_pos), Inches(1.4), Inches(0.9))
        stat_frame = stat_box.text_frame
        stat_frame.text = f"{stat_num}\n{stat_label}"
        
        # Number styling
        num_para = stat_frame.paragraphs[0]
        num_para.alignment = PP_ALIGN.CENTER
        num_run = num_para.runs[0]
        num_run.font.size = Pt(28)
        num_run.font.bold = True
        num_run.font.color.rgb = RGBColor(107, 92, 231)
        num_run.font.name = 'Poppins'
        
        # Label styling
        if len(stat_frame.paragraphs) > 1:
            label_para = stat_frame.paragraphs[1]
        else:
            label_para = stat_frame.add_paragraph()
        label_para.text = stat_label
        label_para.alignment = PP_ALIGN.CENTER
        label_para.font.size = Pt(10)
        label_para.font.color.rgb = RGBColor(100, 116, 139)
        label_para.font.name = 'Inter'
        
        # Box styling
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Save presentation
    prs.save('CodeXBit_Proposal.pptx')
    print("✓ PowerPoint presentation created successfully: CodeXBit_Proposal.pptx")

if __name__ == "__main__":
    create_presentation()

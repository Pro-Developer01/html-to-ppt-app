#!/usr/bin/env python3
"""
Script to convert HTML slides to PowerPoint using screenshot approach
This preserves the exact visual appearance of the HTML
"""

import os
import subprocess
import sys
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

def html_to_image(html_file, output_image, width=1280, height=720):
    """Convert HTML file to PNG image using headless browser"""
    try:
        # Try using playwright (better quality)
        from playwright.sync_api import sync_playwright
        
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={'width': width, 'height': height})
            page.goto(f'file://{os.path.abspath(html_file)}')
            page.wait_for_timeout(1000)  # Wait for rendering
            page.screenshot(path=output_image, full_page=False)
            browser.close()
            return True
    except ImportError:
        print("Playwright not available, trying selenium...")
        
    try:
        # Fallback to selenium
        from selenium import webdriver
        from selenium.webdriver.chrome.options import Options
        
        chrome_options = Options()
        chrome_options.add_argument('--headless')
        chrome_options.add_argument('--no-sandbox')
        chrome_options.add_argument('--disable-dev-shm-usage')
        chrome_options.add_argument(f'--window-size={width},{height}')
        
        driver = webdriver.Chrome(options=chrome_options)
        driver.get(f'file://{os.path.abspath(html_file)}')
        driver.save_screenshot(output_image)
        driver.quit()
        return True
    except Exception as e:
        print(f"Selenium failed: {e}")
        
    return False

def create_html_files():
    """Create proper HTML files from txt files"""
    # Read slide1.txt and save as slide1.html
    with open('slide1.txt', 'r') as f:
        content = f.read()
    with open('slide1.html', 'w') as f:
        f.write(content)
    
    # Read slide2.txt and save as slide2.html
    with open('slide2.txt', 'r') as f:
        content = f.read()
    with open('slide2.html', 'w') as f:
        f.write(content)
    
    print("✓ Created HTML files")

def create_presentation_from_images():
    """Create PowerPoint from screenshot images"""
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Add slide 1
    if os.path.exists('slide1.png'):
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        slide1.shapes.add_picture('slide1.png', 0, 0, width=prs.slide_width, height=prs.slide_height)
        print("✓ Added slide 1")
    
    # Add slide 2
    if os.path.exists('slide2.png'):
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        slide2.shapes.add_picture('slide2.png', 0, 0, width=prs.slide_width, height=prs.slide_height)
        print("✓ Added slide 2")
    
    # Save presentation
    prs.save('CodeXBit_Proposal.pptx')
    print("✓ PowerPoint presentation created: CodeXBit_Proposal.pptx")

def main():
    print("Converting HTML slides to PowerPoint...\n")
    
    # Step 1: Create HTML files
    create_html_files()
    
    # Step 2: Convert HTML to images
    print("\nConverting HTML to images...")
    success1 = html_to_image('slide1.html', 'slide1.png')
    success2 = html_to_image('slide2.html', 'slide2.png')
    
    if not (success1 and success2):
        print("\n⚠ Could not convert HTML to images automatically.")
        print("Please install playwright: pip install playwright && playwright install chromium")
        print("Or install selenium: pip install selenium")
        return
    
    # Step 3: Create PowerPoint from images
    print("\nCreating PowerPoint presentation...")
    create_presentation_from_images()
    
    print("\n✅ Done! Open CodeXBit_Proposal.pptx to view your presentation.")

if __name__ == "__main__":
    main()
